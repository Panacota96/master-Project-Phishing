#!/usr/bin/env python3
"""Generate management_deck.pptx from management_deck.md using python-pptx."""

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError as exc:
    raise SystemExit("python-pptx is required. Install with: python -m pip install python-pptx") from exc


def parse_markdown_slides(md_text):
    slides = []
    current = []
    for line in md_text.splitlines():
        if line.strip() == "---":
            if current:
                slides.append("\n".join(current).strip())
                current = []
            continue
        current.append(line)
    if current:
        slides.append("\n".join(current).strip())
    return [s for s in slides if s]


def build_pptx(slides, output_path):
    prs = Presentation()
    title_only_layout = prs.slide_layouts[5]
    title_and_content_layout = prs.slide_layouts[1]

    for slide_md in slides:
        lines = [l for l in slide_md.splitlines() if l.strip()]
        if not lines:
            continue
        title = lines[0].lstrip("# ").strip()
        body = "\n".join(line.lstrip("# ") for line in lines[1:]).strip()

        layout = title_only_layout if not body else title_and_content_layout
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if body:
            slide.placeholders[1].text = body

    prs.save(output_path)


def main():
    md_path = Path("management_deck.md")
    out_path = Path("management_deck.pptx")
    if not md_path.exists():
        raise SystemExit("management_deck.md not found")
    slides = parse_markdown_slides(md_path.read_text(encoding="utf-8"))
    build_pptx(slides, out_path)
    print(f"Generated {out_path}")


if __name__ == "__main__":
    main()
